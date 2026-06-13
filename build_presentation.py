from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── COLOURS ────────────────────────────────────────────────────────────────────
SAGE       = RGBColor(0x3D, 0x6B, 0x52)
SAGE_LT    = RGBColor(0xEB, 0xF2, 0xEE)
GOLD       = RGBColor(0xA0, 0x85, 0x58)
GOLD_LT    = RGBColor(0xC9, 0xAA, 0x7C)
CREAM      = RGBColor(0xF5, 0xF0, 0xE8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1A, 0x16, 0x14)
INK_SOFT   = RGBColor(0x4A, 0x44, 0x38)
INK_MUTE   = RGBColor(0x8A, 0x80, 0x70)
RED        = RGBColor(0xB0, 0x3A, 0x2A)
RED_LT     = RGBColor(0xFA, 0xEA, 0xE7)
AMBER      = RGBColor(0xB0, 0x6A, 0x1A)
AMBER_LT   = RGBColor(0xFE, 0xF3, 0xE2)
DARK_BG    = RGBColor(0x1A, 0x16, 0x14)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── HELPERS ────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = Pt(0)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    return shape

def add_text(slide, text, l, t, w, h,
             font_name="Calibri", size=Pt(12), bold=False, italic=False,
             color=INK_SOFT, align=PP_ALIGN.LEFT, wrap=True,
             v_anchor=None):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, l, t, w, h,
                  font_name="Calibri", size=Pt(11), bold=False,
                  color=INK_SOFT, align=PP_ALIGN.LEFT, spacing=Pt(6)):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
    return txb

def add_bullet_para(tf, text, font_name="Calibri", size=Pt(11),
                    color=INK_SOFT, indent=False):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(5)
    if indent:
        p.level = 1
    run = p.add_run()
    run.text = ("    •  " if indent else "•  ") + text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    return p

def slide_number(slide, n, total=7):
    add_text(slide, f"{n} / {total}",
             11.8, 7.1, 1.3, 0.3,
             size=Pt(9), color=INK_MUTE, align=PP_ALIGN.RIGHT)

def footer_bar(slide, text="HealthLedgerAI  ·  healthledgerai.com  ·  Confidential"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=DARK_BG)
    add_text(slide, text,
             0.4, 7.12, 12, 0.3,
             size=Pt(9), color=RGBColor(0xAA, 0xA0, 0x90),
             align=PP_ALIGN.LEFT)

def accent_bar(slide, color=SAGE):
    add_rect(slide, 0, 0, 0.18, 7.1, fill=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

add_rect(s, 0, 0, 13.33, 7.5, fill=DARK_BG)                     # dark background
add_rect(s, 0, 0, 0.35, 7.5, fill=SAGE)                          # left accent
add_rect(s, 0.35, 5.6, 12.98, 0.04, fill=GOLD)                   # gold divider

# tag
add_text(s, "CONFIDENTIAL  ·  RESEARCH PROTOTYPE  ·  JUNE 2026",
         1.0, 1.2, 11, 0.4,
         size=Pt(9), color=GOLD, bold=False)

# main title
add_text(s, "Billing Compliance\nfor Private Clinics",
         1.0, 1.75, 11, 2.0,
         font_name="Georgia", size=Pt(48), bold=False,
         color=WHITE)

# subtitle
add_text(s, "A research prototype that automatically detects billing errors\nin private healthcare — tested on 506 real invoices",
         1.0, 3.85, 9, 0.9,
         size=Pt(14), color=RGBColor(0xCC, 0xC5, 0xB8), italic=True)

# presenter block
add_rect(s, 1.0, 5.0, 5.5, 1.1, fill=RGBColor(0x2A, 0x26, 0x22))
add_text(s, "Maria Polychroniadou",
         1.2, 5.08, 5, 0.35,
         font_name="Georgia", size=Pt(15), color=WHITE)
add_text(s, "PhD Researcher  ·  HealthLedgerAI Founder\nBSc International Management  ·  MSc Accounting & Finance",
         1.2, 5.42, 5.2, 0.6,
         size=Pt(10), color=RGBColor(0xAA, 0xA0, 0x90))

# presented to
add_text(s, "Presented to",
         7.2, 5.08, 4, 0.3,
         size=Pt(9), color=GOLD)
add_text(s, "Embryolab Fertility Clinic\nThessaloniki, Greece",
         7.2, 5.38, 5, 0.65,
         font_name="Georgia", size=Pt(16), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
accent_bar(s)
footer_bar(s)
slide_number(s, 2)

add_rect(s, 0.18, 0, 13.15, 0.7, fill=CREAM)
add_text(s, "THE PROBLEM", 0.5, 0.12, 4, 0.35,
         size=Pt(9), color=GOLD, bold=True)

add_text(s, "Private clinics lose money on billing errors they cannot see",
         0.5, 0.82, 12.5, 1.0,
         font_name="Georgia", size=Pt(30), color=INK)

add_rect(s, 0.5, 1.9, 12.33, 0.03, fill=GOLD_LT)

# Four problem boxes
boxes = [
    ("🔴  Duplicate Invoices",
     "The same procedure billed twice to the same insurer — often months apart. Insurers detect these and claw back payments with penalties."),
    ("🔴  Upcoding",
     "A procedure billed above its EOPYY or insurer maximum rate. A single audit trigger that puts your entire insurer contract at risk."),
    ("🔴  Phantom Billing",
     "An invoice raised for a date when the patient has no recorded appointment. In a routine audit, this looks like fraud."),
    ("🟡  Unbundling",
     "Procedure components that should be billed as one bundled code are billed individually — inflating the total without clinical justification."),
]

bx = [0.5, 3.75, 7.0, 9.8]
by = 2.1
bw = 2.9
bh = 1.6

for i, (title, body) in enumerate(boxes):
    x = bx[i % 4] if i < 4 else bx[i % 4]
    add_rect(s, bx[i], by, bw, bh,
             fill=RED_LT if "🔴" in title else AMBER_LT,
             line_color=RED if "🔴" in title else AMBER,
             line_width=Pt(0.75))
    add_text(s, title, bx[i]+0.15, by+0.12, bw-0.2, 0.4,
             size=Pt(11), bold=True, color=RED if "🔴" in title else AMBER)
    add_text(s, body, bx[i]+0.15, by+0.52, bw-0.25, 1.0,
             size=Pt(10), color=INK_SOFT)

add_text(s, "In a dataset of 506 invoices modelled on private clinic billing, HealthLedgerAI detected 70 alerts across 8 error types.",
         0.5, 5.85, 12, 0.4,
         size=Pt(11), italic=True, color=INK_MUTE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
accent_bar(s)
footer_bar(s)
slide_number(s, 3)

add_rect(s, 0.18, 0, 13.15, 0.7, fill=CREAM)
add_text(s, "THE SOLUTION", 0.5, 0.12, 4, 0.35,
         size=Pt(9), color=GOLD, bold=True)

add_text(s, "HealthLedgerAI",
         0.5, 0.82, 8, 0.75,
         font_name="Georgia", size=Pt(36), color=INK)
add_text(s, "Automated billing compliance — in seconds, not weeks",
         0.5, 1.57, 9, 0.4,
         size=Pt(14), italic=True, color=INK_SOFT)
add_rect(s, 0.5, 2.05, 12.33, 0.03, fill=GOLD_LT)

# Three stat boxes
stats = [("11", "Detectors running\nin parallel"),
         ("506", "Invoices analysed\nin the test"),
         ("70", "Alerts flagged\nautomatically")]
sx = [0.5, 4.85, 9.2]
for i, (num, lbl) in enumerate(stats):
    add_rect(s, sx[i], 2.2, 3.7, 1.55, fill=SAGE_LT, line_color=SAGE, line_width=Pt(0.75))
    add_text(s, num, sx[i], 2.28, 3.7, 0.85,
             font_name="Georgia", size=Pt(52), color=SAGE, align=PP_ALIGN.CENTER)
    add_text(s, lbl, sx[i], 3.05, 3.7, 0.55,
             size=Pt(10), color=INK_SOFT, align=PP_ALIGN.CENTER)

add_text(s, "How it works",
         0.5, 3.95, 6, 0.35,
         font_name="Georgia", size=Pt(16), color=INK)

steps = [
    "01   Upload your billing data  —  invoices, appointments, procedure rates, insurer networks",
    "02   The prototype runs 11 detectors automatically across the full dataset",
    "03   Output: a plain Excel report with every alert, its risk level, and a recommended action",
    "04   Your billing team works through the report  —  no special software needed",
]
for i, step in enumerate(steps):
    add_rect(s, 0.5, 4.4 + i*0.55, 12.33, 0.48,
             fill=CREAM if i % 2 == 0 else WHITE,
             line_color=RGBColor(0xDD, 0xD5, 0xC5), line_width=Pt(0.5))
    add_text(s, step, 0.65, 4.43 + i*0.55, 12, 0.42,
             size=Pt(11), color=INK_SOFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DEMO OUTPUT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
accent_bar(s)
footer_bar(s)
slide_number(s, 4)

add_rect(s, 0.18, 0, 13.15, 0.7, fill=CREAM)
add_text(s, "LIVE DEMO RESULTS", 0.5, 0.12, 4, 0.35,
         size=Pt(9), color=GOLD, bold=True)

add_text(s, "What the alert report looks like",
         0.5, 0.82, 9, 0.65,
         font_name="Georgia", size=Pt(30), color=INK)
add_rect(s, 0.5, 1.52, 12.33, 0.03, fill=GOLD_LT)

# Summary bar
add_rect(s, 0.5, 1.65, 12.33, 0.55, fill=DARK_BG)
add_text(s, "506 invoices analysed     |     70 alerts flagged     |     53 HIGH priority     |     17 MEDIUM priority",
         0.7, 1.72, 12, 0.38,
         size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER, bold=True)

alerts = [
    ("HIGH", "🔴", "D1 · Duplicate Invoice",
     "Vendor 'Radiology Dept'  |  Amount €145  |  Date 2026-02-03",
     "Same invoice billed twice. Reject duplicate before payment to avoid insurer clawback."),
    ("HIGH", "🔴", "D2 · Upcoding",
     "Procedure billed at €806  |  Maximum allowed: €620  |  Excess: €186",
     "Withhold payment. Request clinical justification. Risk of EOPYY audit flag."),
    ("HIGH", "🔴", "D4 · Phantom Billing",
     "Patient P047  |  Invoice date: 2026-03-14  |  No appointment record found",
     "Verify patient was present. Do not submit to insurer without confirmation."),
    ("MEDIUM", "🟡", "D11 · Repeated Same-Day Procedure",
     "Same procedure billed 3× on the same day for the same patient",
     "Verify clinical notes confirm multiple administrations before submission."),
]

ay = 2.35
for risk, icon, det, finding, action in alerts:
    bg = RED_LT if risk == "HIGH" else AMBER_LT
    bc = RED if risk == "HIGH" else AMBER
    add_rect(s, 0.5, ay, 12.33, 1.0, fill=bg, line_color=bc, line_width=Pt(1.0))
    add_rect(s, 0.5, ay, 0.08, 1.0, fill=bc)
    add_text(s, f"{icon}  {det}  —  {risk} RISK",
             0.72, ay + 0.06, 11.8, 0.32,
             size=Pt(10), bold=True, color=bc)
    add_text(s, finding,
             0.72, ay + 0.36, 11.8, 0.28,
             size=Pt(10), color=INK, bold=False)
    add_text(s, f"Action: {action}",
             0.72, ay + 0.64, 11.8, 0.28,
             size=Pt(9.5), italic=True, color=INK_MUTE)
    ay += 1.08

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WHY EMBRYOLAB
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
accent_bar(s)
footer_bar(s)
slide_number(s, 5)

add_rect(s, 0.18, 0, 13.15, 0.7, fill=CREAM)
add_text(s, "RELEVANCE TO EMBRYOLAB", 0.5, 0.12, 5, 0.35,
         size=Pt(9), color=GOLD, bold=True)

add_text(s, "Why fertility clinics are a natural fit",
         0.5, 0.82, 10, 0.65,
         font_name="Georgia", size=Pt(30), color=INK)
add_rect(s, 0.5, 1.52, 12.33, 0.03, fill=GOLD_LT)

points = [
    ("EOPYY + Private Insurers Simultaneously",
     "IVF is covered by EOPYY up to 3 cycles for women under 50. Clinics bill two systems with different rules for the same patient — mixed billing is where errors concentrate."),
    ("High Procedure Values",
     "A single IVF cycle costs €3,000–€6,000. One upcoding or duplicate billing error at this scale represents significant financial exposure per incident."),
    ("Multi-Step Billing Sequences",
     "Each cycle generates 8–12 separate billing events: consultations, lab work, anaesthesia, embryology, medications. Date mismatches and duplicates are common at this volume."),
    ("ISO-Certified Environment",
     "Embryolab holds ISO 9001:2015 and EN 15224 certification — meaning it already operates to documented process standards. A billing compliance tool fits naturally into that framework."),
]

py = 1.7
for title, body in points:
    add_rect(s, 0.5, py, 0.06, 0.9, fill=SAGE)
    add_text(s, title, 0.7, py + 0.02, 12, 0.3,
             size=Pt(11), bold=True, color=INK)
    add_text(s, body, 0.7, py + 0.33, 12, 0.55,
             size=Pt(10.5), color=INK_SOFT)
    add_rect(s, 0.5, py + 0.95, 12.33, 0.02, fill=RGBColor(0xDD, 0xD5, 0xC5))
    py += 1.08

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — THE ASK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
accent_bar(s)
footer_bar(s)
slide_number(s, 6)

add_rect(s, 0.18, 0, 13.15, 0.7, fill=CREAM)
add_text(s, "WHAT I AM ASKING FOR", 0.5, 0.12, 5, 0.35,
         size=Pt(9), color=GOLD, bold=True)

add_text(s, "20 minutes and your honest opinion",
         0.5, 0.82, 10, 0.65,
         font_name="Georgia", size=Pt(30), color=INK)
add_text(s, "Δεν πουλάω τίποτα — ζητώ εμπειρογνωμοσύνη",
         0.5, 1.48, 10, 0.38,
         size=Pt(13), italic=True, color=INK_MUTE)
add_rect(s, 0.5, 1.92, 12.33, 0.03, fill=GOLD_LT)

# context box
add_rect(s, 0.5, 2.05, 12.33, 0.95, fill=CREAM, line_color=GOLD_LT, line_width=Pt(0.5))
add_text(s, "I am a doctoral researcher, not a salesperson. HealthLedgerAI is a research prototype I am developing as part of my PhD on billing compliance in private healthcare in Greece and Cyprus. I signed a Non-Disclosure Agreement with you before this meeting — everything I share stays protected, and everything you share with me stays within the bounds of that agreement.",
         0.7, 2.1, 12, 0.85,
         size=Pt(11), color=INK_SOFT)

# four asks
ask_items = [
    ("01", "Do the billing error types this prototype detects reflect real problems you encounter in your billing process?"),
    ("02", "Are there error types I have missed that are common in fertility clinic billing in Greece?"),
    ("03", "Would a tool like this be useful to your billing or administration team — even at prototype stage?"),
    ("04", "Would you provide a short written statement that you reviewed this prototype? (Voluntary — for my doctoral research record only.)"),
]

add_rect(s, 0.5, 3.1, 12.33, 0.38, fill=SAGE)
add_text(s, "Four questions I would like to ask you:",
         0.7, 3.15, 12, 0.3,
         size=Pt(12), bold=True, color=WHITE)

iy = 3.55
for num, text in ask_items:
    add_rect(s, 0.5, iy, 0.7, 0.7, fill=SAGE_LT)
    add_text(s, num, 0.5, iy + 0.12, 0.7, 0.4,
             font_name="Georgia", size=Pt(18), color=SAGE, align=PP_ALIGN.CENTER)
    add_rect(s, 1.25, iy, 11.58, 0.7,
             fill=WHITE if int(num) % 2 == 1 else CREAM,
             line_color=RGBColor(0xDD, 0xD5, 0xC5), line_width=Pt(0.5))
    add_text(s, text, 1.4, iy + 0.12, 11.2, 0.55,
             size=Pt(11), color=INK_SOFT)
    iy += 0.76

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ABOUT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=DARK_BG)
add_rect(s, 0, 0, 0.35, 7.5, fill=SAGE)
add_rect(s, 0.35, 5.55, 12.98, 0.04, fill=GOLD)
footer_bar(s, "healthledgerai.com  ·  info@healthledgerai.com  ·  All content covered by NDA")
slide_number(s, 7)

add_text(s, "ABOUT THE RESEARCHER", 0.7, 0.35, 5, 0.35,
         size=Pt(9), color=GOLD, bold=True)
add_text(s, "Maria Polychroniadou",
         0.7, 0.78, 9, 0.75,
         font_name="Georgia", size=Pt(38), color=WHITE)
add_text(s, "PhD Researcher  ·  HealthLedgerAI Founder",
         0.7, 1.52, 8, 0.38,
         size=Pt(14), italic=True, color=RGBColor(0xCC, 0xC5, 0xB8))

# left column — academic
add_rect(s, 0.7, 2.1, 5.6, 0.42, fill=SAGE)
add_text(s, "Academic Background",
         0.9, 2.16, 5.2, 0.32,
         size=Pt(11), bold=True, color=WHITE)

acad = [
    "BSc  International Management",
    "MSc  Accounting & Finance",
    "PhD candidate — AI & billing compliance in private healthcare (Greece & Cyprus)",
    "Erasmus for Young Entrepreneurs — Thessaloniki, 2026",
]
ay2 = 2.6
for item in acad:
    add_rect(s, 0.7, ay2, 0.06, 0.42, fill=GOLD)
    add_text(s, item, 0.85, ay2 + 0.04, 5.3, 0.38,
             size=Pt(11), color=RGBColor(0xCC, 0xC5, 0xB8))
    ay2 += 0.5

# right column — project
add_rect(s, 7.1, 2.1, 5.8, 0.42, fill=GOLD)
add_text(s, "HealthLedgerAI",
         7.3, 2.16, 5.4, 0.32,
         size=Pt(11), bold=True, color=DARK_BG)

proj = [
    "Pre-launch research prototype",
    "11 billing error detectors built and tested",
    "Applied to EIT Jumpstarter 2026",
    "Target: EIC Pre-Accelerator 2027 (€500K–€1M grant)",
]
py2 = 2.6
for item in proj:
    add_rect(s, 7.1, py2, 0.06, 0.42, fill=SAGE)
    add_text(s, item, 7.25, py2 + 0.04, 5.5, 0.38,
             size=Pt(11), color=RGBColor(0xCC, 0xC5, 0xB8))
    py2 += 0.5

add_text(s, "healthledgerai.com  ·  info@healthledgerai.com",
         0.7, 5.72, 12, 0.38,
         size=Pt(11), color=RGBColor(0x8A, 0x80, 0x70), align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────────
prs.save("PRESENTATION-EMBRYOLAB.pptx")
print("Saved: PRESENTATION-EMBRYOLAB.pptx")
